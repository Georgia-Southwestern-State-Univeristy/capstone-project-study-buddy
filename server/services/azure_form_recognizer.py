""" Helper functions for extracting text from various file types using Azure Form Recognizer """

"""Step 1: Import necessary modules"""
import io
import os
import logging
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential

# ---- Computer Vision for fallback OCR ----
from azure.cognitiveservices.vision.computervision import ComputerVisionClient
from azure.cognitiveservices.vision.computervision.models import OperationStatusCodes
from msrest.authentication import CognitiveServicesCredentials
import time

# Import additional libraries
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation

"""Step 2: Define the helper functions"""
def get_form_recognizer_client():
    """
    Returns the Azure Form Recognizer client, using environment variables:
    - AZURE_FORM_RECOGNIZER_ENDPOINT
    - AZURE_FORM_RECOGNIZER_KEY
    """
    endpoint = os.getenv("AZURE_FORM_RECOGNIZER_ENDPOINT")
    key = os.getenv("AZURE_FORM_RECOGNIZER_KEY")
    if not endpoint or not key:
        logging.warning("Form Recognizer endpoint/key not found in env variables.")
    return DocumentAnalysisClient(endpoint, AzureKeyCredential(key))


def get_computer_vision_client():
    """
    Returns the Azure Computer Vision client, using environment variables:
    - AZURE_COMPUTER_VISION_ENDPOINT
    - AZURE_COMPUTER_VISION_KEY
    """
    endpoint = os.getenv("AZURE_COMPUTER_VISION_ENDPOINT")
    key = os.getenv("AZURE_COMPUTER_VISION_KEY")
    if not endpoint or not key:
        logging.warning("Computer Vision endpoint/key not found in env variables.")
    return ComputerVisionClient(endpoint, CognitiveServicesCredentials(key))


def extract_text_from_file(file_content: bytes, file_mime_type: str) -> str:
    """
    Main entry point for text extraction. 
    1) Dispatches .docx/.xlsx/.pptx to dedicated parsers.
    2) For PDF/images, tries Form Recognizer's 'prebuilt-read' model.
    3) Optionally, if we detect an image and want more robust OCR, fallback to Computer Vision.
    """
    try:
        if file_mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            # .docx
            logging.info("Processing a .docx file")
            return extract_text_from_docx(file_content)

        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            # .xlsx
            logging.info("Processing an .xlsx file")
            return extract_text_from_xlsx(file_content)

        elif file_mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            # .pptx
            logging.info("Processing a .pptx file")
            return extract_text_from_pptx(file_content)

        elif file_mime_type == 'text/plain':
            # .txt
            logging.info("Processing a .txt file")
            return extract_text_from_txt(file_content)

        else:
            # If it's PDF or an image, attempt Form Recognizer's prebuilt-read first
            logging.info("Trying Form Recognizer for PDF/image")
            extracted_text = try_form_recognizer(file_content, file_mime_type)
            if extracted_text.strip():
                return extracted_text
            else:
                # If no text came out, we can fallback to Computer Vision for images
                if file_mime_type.startswith('image/'):
                    logging.info("Form Recognizer returned no text; trying Computer Vision OCR fallback.")
                    extracted_text_cv = try_computer_vision_ocr(file_content)
                    return extracted_text_cv
                return extracted_text  # For PDF with no text, just return empty string

    except Exception as e:
        logging.error(f"Error extracting text from file: {e}", exc_info=True)
        return ""


def try_form_recognizer(file_content: bytes, file_mime_type: str) -> str:
    """
    Uses Azure Form Recognizer's 'prebuilt-read' model to extract text.
    The crucial part is passing content_type=file_mime_type so it can handle images.
    """
    try:
        client = get_form_recognizer_client()
        if not client:
            logging.error("Form Recognizer client is not available. Check environment variables.")
            return ""

        poller = client.begin_analyze_document(
            "prebuilt-read",
            document=file_content,
            content_type=file_mime_type  # <--- Pass MIME type
        )
        result = poller.result()

        extracted_text = ""
        for page in result.pages:
            for line in page.lines:
                extracted_text += line.content + "\n"
        return extracted_text

    except Exception as e:
        logging.error(f"Form Recognizer failed: {e}", exc_info=True)
        return ""


def try_computer_vision_ocr(image_bytes: bytes) -> str:
    """
    Fallback to Azure Computer Vision OCR if needed.
    Uses the 'Read' API (v3) in a synchronous polling manner.
    """
    extracted_text = ""
    try:
        cv_client = get_computer_vision_client()
        if not cv_client:
            logging.error("Computer Vision client is not available. Check environment variables.")
            return ""

        # Initiate the read operation (raw=True returns the HTTP response directly)
        read_response = cv_client.read_in_stream(io.BytesIO(image_bytes), raw=True)
        operation_location = read_response.headers["Operation-Location"]
        operation_id = operation_location.split("/")[-1]

        # Poll for the result
        while True:
            read_result = cv_client.get_read_result(operation_id)
            if read_result.status not in [OperationStatusCodes.running, OperationStatusCodes.not_started]:
                break
            time.sleep(1)

        # If the operation is successful, parse out text
        if read_result.status == OperationStatusCodes.succeeded:
            for text_result in read_result.analyze_result.read_results:
                for line in text_result.lines:
                    extracted_text += line.text + "\n"

    except Exception as e:
        logging.error(f"Computer Vision OCR failed: {e}", exc_info=True)

    return extracted_text


def extract_text_from_docx(file_content: bytes) -> str:
    """
    Extract text from .docx using python-docx
    """
    extracted_text = ""
    try:
        document = DocxDocument(io.BytesIO(file_content))
        for para in document.paragraphs:
            extracted_text += para.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .docx file: {e}", exc_info=True)
    return extracted_text


def extract_text_from_xlsx(file_content: bytes) -> str:
    """
    Extract text from .xlsx using openpyxl
    """
    extracted_text = ""
    try:
        workbook = openpyxl.load_workbook(io.BytesIO(file_content), data_only=True)
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = [str(cell) if cell is not None else "" for cell in row]
                extracted_text += "\t".join(row_text) + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .xlsx file: {e}", exc_info=True)
    return extracted_text


def extract_text_from_pptx(file_content: bytes) -> str:
    """
    Extract text from .pptx using python-pptx
    """
    extracted_text = ""
    try:
        presentation = Presentation(io.BytesIO(file_content))
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += shape.text + "\n"
    except Exception as e:
        logging.error(f"Error extracting text from .pptx file: {e}", exc_info=True)
    return extracted_text


def extract_text_from_txt(file_content: bytes) -> str:
    """
    Extract text from a simple text file
    """
    return file_content.decode('utf-8', errors='ignore')


# Define a function to check if a file type is supported
ALLOWED_MIME_TYPES = [
            'application/pdf',
            'image/jpeg',
            'image/png',
            'image/tiff',
            'image/bmp',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # .docx
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',        # .xlsx
            'application/vnd.openxmlformats-officedocument.presentationml.presentation', # .pptx
            'application/vnd.ms-powerpoint',
            'text/plain',  # .txt
            # Add more MIME types as needed
        ]